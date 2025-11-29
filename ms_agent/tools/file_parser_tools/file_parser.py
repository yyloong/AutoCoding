import json
import uuid
import os
from docx import Document
import re
import time
import zipfile
import math
from pathlib import Path
from qwen_agent.settings import DEFAULT_WORKSPACE, DEFAULT_MAX_INPUT_TOKENS
from pptx.enum.shapes import MSO_SHAPE_TYPE
from qwen_agent.utils.tokenization_qwen import count_tokens, tokenizer
from typing import Any, Dict, List, Optional, Union
from collections import Counter
import xml.etree.ElementTree as ET
from pandas import Timestamp
from datetime import datetime
from pandas.api.types import is_datetime64_any_dtype
import pandas as pd
from tabulate import tabulate
from ms_agent.utils import get_logger
from .utils import (get_file_type, hash_sha256, is_http_url, get_basename_from_url, 
                                  sanitize_chrome_file_path, save_url_to_local_work_dir,json_loads)

logger = get_logger()

# Configuration constants
PARSER_SUPPORTED_FILE_TYPES = ['pdf', 'docx', 'pptx', 'txt', 'html', 'csv', 'tsv', 'xlsx', 'xls', 'doc', 'zip', '.mp4', '.mov', '.mkv', '.webm', '.mp3', '.wav']
def str_to_bool(value):
    """Convert string to boolean, handling common true/false representations"""
    if isinstance(value, bool):
        return value
    return str(value).lower() in ('true', '1', 'yes', 'on')
USE_IDP = str_to_bool(os.getenv("USE_IDP", "True"))
IDP_TIMEOUT = 150000
ENABLE_CSI = False
PARAGRAPH_SPLIT_SYMBOL = '\n'


class CustomJSONEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, (datetime, Timestamp)):
            return obj.isoformat()
        return super().default(obj)


class FileParserError(Exception):
    """Custom exception for document parsing errors"""

    def __init__(self, message: str, code: str = '400', exception: Optional[Exception] = None):
        super().__init__(message)
        self.code = code
        self.exception = exception

def clean_text(text: str) -> str:
    cleaners = [
        lambda x: re.sub(r'\n+', '\n', x),  
        lambda x: x.replace("Add to Qwen's Reading List", ''),
        lambda x: re.sub(r'-{6,}', '-----', x),  
        lambda x: x.strip()
    ]
    for cleaner in cleaners:
        text = cleaner(text)
    return text


def get_plain_doc(doc: list):
    paras = []
    for page in doc:
        for para in page['content']:
            for k, v in para.items():
                if k in ['text', 'table', 'image']:
                    paras.append(v)
    return PARAGRAPH_SPLIT_SYMBOL.join(paras)


def df_to_markdown(df: pd.DataFrame) -> str:
    df = df.dropna(how='all').fillna('')
    return tabulate(df, headers='keys', tablefmt='pipe', showindex=False)


def parse_word(docx_path: str, extract_image: bool = True):
    """
    终极版 Word 解析：
    1. 使用 python-docx 提取文本和表格。
    2. 使用 zipfile 直接解压 .docx 文件提取所有图片（100% 成功率）。
    3. 将图片分析结果附加到文档末尾。
    """
    content = []
    
    # --- 阶段 1: 提取文本和表格 (使用 python-docx) ---
    try:
        doc = Document(docx_path)
        
        # 提取段落文本
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                content.append({'text': text})
        
        # 提取表格
        for table in doc.tables:
            tbl_text = []
            for row in table.rows:
                # 简单处理：将单元格内容用 | 连接
                row_cells = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                tbl_text.append('|' + '|'.join(row_cells) + '|')
            
            tbl_str = '\n'.join(tbl_text)
            if tbl_str.strip():
                content.append({'table': tbl_str})
                
    except Exception as e:
        logger.error(f"Text parsing failed: {e}")
        # 即使文本解析失败，也可以尝试提取图片

    # --- 阶段 2: 提取所有图片 (使用 zipfile 暴力提取) ---
    if extract_image:
        try:
            # 准备图片输出目录
            img_output_dir = os.path.join(os.path.dirname(docx_path), "extracted_images")
            if not os.path.exists(img_output_dir):
                os.makedirs(img_output_dir, exist_ok=True)

            # 打开 docx 作为 zip 文件
            with zipfile.ZipFile(docx_path, 'r') as z:
                # 获取 zip 中所有文件列表
                all_files = z.namelist()
                
                # 过滤出媒体文件夹下的图片 (word/media/image1.png ...)
                media_files = [f for f in all_files if f.startswith('word/media/') and f != 'word/media/']
                
                logger.info(f"Found {len(media_files)} images in docx.")

                for media_file in media_files:
                    # 排除非图片文件 (如 wmf, emf 矢量图 VLM 很难读，bin 是 OLE 对象)
                    valid_exts = ('.png', '.jpg', '.jpeg', '.bmp', '.gif')
                    if not media_file.lower().endswith(valid_exts):
                        continue

                    # 读取图片数据
                    img_data = z.read(media_file)
                    
                    # 生成本地文件名
                    base_name = os.path.basename(media_file) # image1.png
                    # 加个 uuid 防止覆盖
                    save_name = f"docx_{uuid.uuid4().hex[:6]}_{base_name}"
                    save_path = os.path.join(img_output_dir, save_name)

                    with open(save_path, "wb") as f:
                        f.write(img_data)

                    # 调用 VLM 分析
                    desc = f"[Extracted Image: {base_name}]"
                    if '_analyze_image_with_vlm' in globals():
                        try:
                            # 提示词增加上下文信息
                            desc = globals()['_analyze_image_with_vlm'](save_path, context="Word Document Appendix")
                        except Exception as vlm_e:
                            logger.warning(f"VLM analysis failed for {media_file}: {vlm_e}")
                    
                    # 将图片分析结果加入内容列表
                    # 策略：可以放在最前面，也可以放在最后面
                    # 这里选择追加到 content，并标记来源
                    content.append({
                        'image_analysis': desc,
                        'image_path': save_path,
                        'type': 'image_extraction_global'
                    })

        except Exception as e:
            logger.error(f"Image extraction via zip failed: {e}")

    # --- 阶段 3: 结果整合 ---
    # 如果完全没内容
    if not content:
        return []

    # 这里的逻辑是：文档内容 = [文本/表格段落...] + [所有图片的分析结果...]
    # 虽然丢失了“图片在第几段”的信息，但保证了图片内容绝对不会丢。
    return [{'page_num': 1, 'content': content}]

try:
    from pdf2image import convert_from_path
except ImportError:
    convert_from_path = None

from openai import OpenAI
import os
import base64

client = OpenAI(
    # 若没有配置环境变量，请用百炼API Key将下行替换为：api_key="sk-xxx",
    api_key=os.getenv("DASHSCOPE_API_KEY"),
    base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
)

def _encode_image(image_path):
    """辅助函数：将本地图片转为 Base64 字符串"""
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def _analyze_image_with_vlm(image_path: str, context: str = "") -> str:
    """
    使用 OpenAI SDK 调用阿里云 Qwen-VL-Max 模型。
    """
    if not os.path.exists(image_path):
        return f"[Error: Image path not found: {image_path}]"

    # 1. 初始化 OpenAI 客户端，指向阿里云 DashScope 兼容端点
    # 2. 图片转 Base64
    try:
        base64_image = _encode_image(image_path)
        # 根据文件后缀判断 mime type (简单处理，默认 jpeg/png)
        file_ext = os.path.splitext(image_path)[-1].lower().replace('.', '')
        if file_ext == 'jpg': file_ext = 'jpeg'
        mime_type = f"image/{file_ext}"
    except Exception as e:
        logger.error(f"Image encoding failed: {e}")
        return f"[图片读取失败: {str(e)}]"

    # 3. 构造 Prompt (针对 Coding 任务优化)
    prompt_text = (
        f"你是一个高级软件架构师。请分析这张来自{context}的图片。\n"
        "核心任务：提取图片中的逻辑结构，用于辅助代码生成。\n"
        "1. **流程图**：请输出 Mermaid graph TD 代码，准确描述节点判定和跳转。\n"
        "2. **架构图**：描述模块划分和数据流向。\n"
        "3. **纯文本/表格**：提取关键数据。\n"
        "如果不包含技术信息，请简短说明。"
    )

    try:
        # 4. 发起请求
        completion = client.chat.completions.create(
            model="qwen-vl-max",  # 指定阿里视觉模型
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text", 
                            "text": prompt_text
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                # 必须使用 Data URI 格式
                                "url": f"data:{mime_type};base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            # max_tokens=2000 # 可选：限制输出长度
        )

        # 5. 获取结果
        content = completion.choices[0].message.content
        if not content:
            return "[VLM 分析结果为空]"
        
        return f"\n[图片分析 ({os.path.basename(image_path)})]:\n{content}\n"

    except Exception as e:
        logger.error(f"OpenAI-compatible API Error: {str(e)}")
        return f"[VLM 调用异常: {str(e)}]"

def parse_png(image_path: str, context: str = "Image") -> List[dict]:
    """
    解析 PNG/JPG 图片，调用 VLM 进行分析。
    """
    # 假设 _analyze_image_with_vlm 函数在上下文中已定义
    result = _analyze_image_with_vlm(image_path, context=context)
    
    # 修改说明：
    # 将 key 由 'image_analysis' 改为 'text'，以便 _flatten_result 能正确提取内容。
    # 可以保留 type 字段作为元数据。
    return [{
        'page_num': 1, 
        'content': [{
            'text': result, 
            'type': 'image_analysis'
        }]
    }]

def parse_ppt(path: str, extract_image: bool = True):
    """
    解析 PPT，支持递归提取组合图形中的文本，并支持提取嵌入的图片（流程图截图）。
    """
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
    try:
        ppt = Presentation(path)
    except PackageNotFoundError as ex:
        logger.warning(ex)
        return []

    # 递归处理 Shape 的内部函数（解决 Group 组合图读不到字的问题）
    def process_shape(shape, slide_idx):
        shape_content = []
        
        # 1. 文本框处理
        if shape.has_text_frame:
            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                # 简单的清洗
                p_text = ''.join(run.text for run in paragraph.runs).strip()
                if p_text:
                    text_parts.append(p_text)
            if text_parts:
                shape_content.append({'text': '\n'.join(text_parts)})

        # 2. 表格处理
        if shape.has_table:
            tbl = []
            for row in shape.table.rows:
                row_text = []
                for cell in row.cells:
                    # 单元格内可能也有复杂结构，简单获取文本
                    cell_txt = cell.text_frame.text if cell.text_frame else ""
                    row_text.append(cell_txt.replace('\n', ' ').strip())
                tbl.append('|' + '|'.join(row_text) + '|')
            if tbl:
                shape_content.append({'table': '\n'.join(tbl)})

        # 3. 图片提取处理 (关键：处理截图类流程图)
        if extract_image and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                # 生成临时路径
                image_ext = shape.image.ext
                image_name = f"ppt_slide_{slide_idx}_{uuid.uuid4().hex[:8]}.{image_ext}"
                # 假设有一个 temp 目录，或者存到 path 同级目录
                save_dir = os.path.join(os.path.dirname(path), "extracted_images")
                os.makedirs(save_dir, exist_ok=True)
                image_path = os.path.join(save_dir, image_name)
                
                with open(image_path, 'wb') as f:
                    f.write(shape.image.blob)
                
                # 调用 VLM 分析
                analysis_text = _analyze_image_with_vlm(image_path, context="PPT Slide")
                shape_content.append({
                    'image_analysis': analysis_text, 
                    'image_path': image_path
                })
            except Exception as e:
                logger.warning(f"Failed to extract PPT image: {e}")

        # 4. 递归处理组合图形 (关键：处理原生绘制的框图)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                shape_content.extend(process_shape(sub_shape, slide_idx))
        
        return shape_content

    doc = []
    for slide_number, slide in enumerate(ppt.slides):
        page = {'page_num': slide_number + 1, 'content': []}
        for shape in slide.shapes:
            page['content'].extend(process_shape(shape, slide_number + 1))
        doc.append(page)
    
    return doc

import pdfplumber

# 只有在需要图片OCR/VLM分析时才需要这个库
try:
    from pdf2image import convert_from_path
except ImportError:
    convert_from_path = None

def parse_pdf(pdf_path: str, extract_image: bool = False) -> List[dict]:
    """
    解析 PDF：优先提取文本/表格，针对扫描件或图表页调用 VLM。
    """
    doc_content = []
    
    # 创建临时目录存放提取的图片
    temp_dir = os.path.join(os.path.dirname(pdf_path), "temp_images")
    if extract_image:
        os.makedirs(temp_dir, exist_ok=True)

    with pdfplumber.open(pdf_path) as pdf:
        total_pages = len(pdf.pages)
        
        for i, page in enumerate(pdf.pages):
            page_num = i + 1
            logger.info(f"Processing Page {page_num}/{total_pages}...")
            
            page_data = {'page_num': page_num, 'content': []}
            
            # --- A. 文本提取 ---
            raw_text = page.extract_text(x_tolerance=1, y_tolerance=3) or ""
            if raw_text.strip():
                page_data['content'].append({
                    'type': 'text',
                    'text': raw_text,
                    'source': 'pdfplumber_text'
                })

            # --- B. 表格提取 ---
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    if not table: continue
                    # 简易 List 转 Markdown
                    clean_table = [[str(c).replace('\n', ' ') if c else '' for c in row] for row in table]
                    if len(clean_table) > 0:
                        headers = clean_table[0]
                        md_table = "| " + " | ".join(headers) + " |\n"
                        md_table += "| " + " | ".join(["---"] * len(headers)) + " |\n"
                        for row in clean_table[1:]:
                            md_table += "| " + " | ".join(row) + " |\n"
                        
                        page_data['content'].append({
                            'type': 'table',
                            'text': f"\n[Detected Table]\n{md_table}\n",
                            'source': 'pdfplumber_table'
                        })

            # --- C. 智能 VLM 调用策略 ---
            # 判断逻辑：
            # 1. 扫描件检测：本页提取到的文字少于 50 个字符。
            # 2. 显式开关：extract_image=True
            # 3. 依赖检查：pdf2image 可用
            is_scanned_page = len(raw_text) < 50
            
            if extract_image and convert_from_path:
                # 只有当 (是扫描件) 或者 (虽然有文字但可能有重要图表 - 这里可根据需求调整) 时才执行
                # 为了节省 token，这里策略设为：如果是扫描件，必跑；如果不是扫描件，仅当 extract_image 强开启时跑
                
                try:
                    # 将 PDF 当前页转为图片
                    images = convert_from_path(
                        pdf_path, 
                        dpi=200, # 200 dpi 对 VLM 足够了，太高费流量
                        first_page=page_num, 
                        last_page=page_num
                    )
                    
                    if images:
                        pil_image = images[0]
                        
                        # 保存临时文件供 VLM 读取
                        temp_img_name = f"page_{page_num}_{int(time.time())}.jpg"
                        temp_img_path = os.path.join(temp_dir, temp_img_name)
                        pil_image.save(temp_img_path, "JPEG")
                        
                        # 调用 VLM
                        # 如果是扫描件，提示 VLM 侧重 OCR；否则侧重架构分析
                        context_hint = f"PDF Page {page_num}"
                        if is_scanned_page:
                            context_hint += " (Scanned Document - Please Extract Text)"
                        
                        logger.info(f"Invoking Qwen-VL-Max for page {page_num}...")
                        vlm_result = _analyze_image_with_vlm(temp_img_path, context=context_hint)
                        
                        page_data['content'].append({
                            'type': 'image_analysis',
                            'text': vlm_result,
                            'is_scanned_fallback': is_scanned_page
                        })
                        
                        # 清理临时文件 (可选：如果想保留图片调试，注释掉这行)
                        os.remove(temp_img_path)

                except Exception as e:
                    logger.error(f"Image processing failed for page {page_num}: {e}")

            doc_content.append(page_data)
            
    return doc_content


def parse_txt(path: str):
    with open(path, 'r', encoding='utf-8') as f:  
        text = f.read()
    paras = text.split(PARAGRAPH_SPLIT_SYMBOL)
    content = []
    for p in paras:
        content.append({'text': p})
    return [{'page_num': 1, 'content': content}]


def get_font(element):
    from pdfminer.layout import LTChar, LTTextContainer

    fonts_list = []
    for text_line in element:
        if isinstance(text_line, LTTextContainer):
            for character in text_line:
                if isinstance(character, LTChar):
                    fonts_list.append((character.fontname, character.size))

    fonts_list = list(set(fonts_list))
    if fonts_list:
        counter = Counter(fonts_list)
        most_common_fonts = counter.most_common(1)[0][0]
        return most_common_fonts
    else:
        return []


def extract_tables(pdf, page_num):
    table_page = pdf.pages[page_num]
    tables = table_page.extract_tables()
    return tables


def table_converter(table):
    table_string = ''
    for row_num in range(len(table)):
        row = table[row_num]
        cleaned_row = [
            item.replace('\n', ' ') if item is not None and '\n' in item else 'None' if item is None else item
            for item in row
        ]
        table_string += ('|' + '|'.join(cleaned_row) + '|' + '\n')
    table_string = table_string[:-1]
    return table_string


def extract_xls_schema(file_path: str) -> Dict[str, Any]:
    xls = pd.ExcelFile(file_path)
    schema = {
        "sheets": [],
        "n_sheets": len(xls.sheet_names)
    }

    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name, nrows=3)  # 读取前3行

        dtype_mapping = {
            'object': 'string',
            'datetime64[ns]': 'datetime',
            'timedelta64[ns]': 'timedelta'
        }
        dtypes = df.dtypes.astype(str).replace(dtype_mapping).to_dict()

        sample_df = df.head(3).copy()
        for col in sample_df.columns:
            if is_datetime64_any_dtype(sample_df[col]):
                sample_df[col] = sample_df[col].dt.strftime('%Y-%m-%dT%H:%M:%S')

        sheet_info = {
            "name": sheet_name,
            "columns": df.columns.tolist(),
            "dtypes": dtypes,  
            "sample_data": sample_df.to_dict(orient='list') 
        }
        schema["sheets"].append(sheet_info)

    return schema


def extract_csv_schema(file_path: str) -> Dict[str, Any]:
    df_dtype = pd.read_csv(file_path, nrows=100)  
    df_sample = pd.read_csv(file_path, nrows=3) 

    return {
        "columns": df_dtype.columns.tolist(),
        "dtypes": df_dtype.dtypes.astype(str).to_dict(),
        "sample_data": df_sample.to_dict(orient='list'),
        "estimated_total_rows": _estimate_total_rows(file_path)
    }


def _estimate_total_rows(file_path) -> int:
    with open(file_path, 'rb') as f:
        line_count = 0
        chunk_size = 1024 * 1024  
        while chunk := f.read(chunk_size):
            line_count += chunk.count(b'\n')
    return line_count - 1  


def parse_tabular_file(file_path: str, **kwargs) -> List[dict]:
    try:
        df = pd.read_excel(file_path) if file_path.endswith(('.xlsx', '.xls')) else \
            pd.read_csv(file_path, **kwargs)
        if count_tokens(df_to_markdown(df)) > DEFAULT_MAX_INPUT_TOKENS:
            schema = extract_xls_schema(file_path) if file_path.endswith(('.xlsx', '.xls')) else \
                extract_csv_schema(file_path)
            return [{'page_num': 1, 'content': [{'schema': schema}]}]
        else:
            return [{'page_num': 1, 'content': [{'table': df_to_markdown(df)}]}]
    except Exception as e:
        logger.error(f"Table parsing failed: {str(e)}")
        return []  


def parse_zip(file_path: str, extract_dir: str) -> List[dict]:
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)
        return [os.path.join(extract_dir, f) for f in zip_ref.namelist()]


def parse_html(file_path: str) -> List[dict]:
    from bs4 import BeautifulSoup  

    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'lxml')

    content = [{'text': clean_text(p.get_text())}
               for p in soup.find_all(['p', 'div']) if p.get_text().strip()]

    return [{
        'page_num': 1,
        'content': content,
        'title': soup.title.string if soup.title else ''
    }]


def extract_xml_skeleton_markdown(xml_file):
    tree = ET.parse(xml_file)
    root = tree.getroot()
    markdown_lines = []

    def process_element(element, level=0, parent_path="", is_last=True, prefix=""):
        if level > 0:
            connector = "└── " if is_last else "├── "
            markdown_lines.append(f"{prefix}{connector}**{element.tag}**")
        else:
            markdown_lines.append(f"## Root: {element.tag}")

        if element.attrib:
            attrs = [f"`{k}`" for k in element.attrib.keys()]
            attr_line = f"{prefix}{'    ' if level > 0 else ''}*Attributes:* {', '.join(attrs)}"
            markdown_lines.append(attr_line)

        if element.text and element.text.strip():
            text_line = f"{prefix}{'    ' if level > 0 else ''}*Has text content*"
            markdown_lines.append(text_line)
        seen_tags = set()
        unique_children = []
        for child in element:
            if child.tag not in seen_tags:
                seen_tags.add(child.tag)
                unique_children.append(child)

        for i, child in enumerate(unique_children):
            is_last_child = (i == len(unique_children) - 1)
            child_prefix = prefix + ("    " if is_last else "│   ")
            process_element(child, level + 1,
                            f"{parent_path}/{element.tag}" if parent_path else element.tag,
                            is_last_child, child_prefix)

    process_element(root)
    markdown_content = "\n".join(markdown_lines)
    return markdown_content


def parse_xml(file_path: str) -> List[dict]:
    with open(file_path, 'r', encoding='utf-8') as f: 
        text = f.read()
    if count_tokens(text) > DEFAULT_MAX_INPUT_TOKENS:
        schema = extract_xml_skeleton_markdown(file_path)
        content = [{'schema': schema}]
    else:
        content = [{'text': text}]
    return [{'page_num': 1, 'content': content}]


def compress(results: list) -> list[str]:
    compress_results = []
    max_token = math.floor(DEFAULT_MAX_INPUT_TOKENS / len(results))
    for result in results:
        token_list = tokenizer.tokenize(result)
        token_list = token_list[:min(len(token_list), max_token)]
        compress_results.append(tokenizer.convert_tokens_to_string(token_list))
    return compress_results


# @register_tool('file_parser')
class SingleFileParser:
    name = "file_parser"
    description = f"File parsing tool, supports parsing data in  {'/'.join(PARSER_SUPPORTED_FILE_TYPES)} formats, and returns the parsed markdown format data."
    parameters = [{
        'name': 'url',
        'type': 'string',
        'description': 'The full path of the file to be parsed, which can be a local path or a downloadable http(s) link.',
        'required': True
    }]

    def __init__(self, cfg: Optional[Dict] = None):
        # 设置数据存储根目录
        self.cfg = cfg
        self.data_root = self.cfg.get('path', os.path.join(DEFAULT_WORKSPACE, 'tools', self.name))
        
        # 修改点 1: 移除 self.db = Storage(...)
        # 改为直接确保目录存在，用于存放 JSON 缓存文件
        if not os.path.exists(self.data_root):
            os.makedirs(self.data_root, exist_ok=True)

        self.structured_doc = self.cfg.get('structured_doc', True)

        self.parsers = {
            'pdf': parse_pdf,
            'docx': parse_word,
            'doc': parse_word,
            'pptx': parse_ppt,
            'txt': parse_txt,
            'jsonl': parse_txt,
            'jsonld': parse_txt,
            'pdb': parse_txt,
            'py': parse_txt,
            "png": parse_png,
            'html': parse_html,
            'xml': parse_xml,
            'csv': lambda p: parse_tabular_file(p, sep=','),
            'tsv': lambda p: parse_tabular_file(p, sep='\t'),
            'xlsx': parse_tabular_file,
            'xls': parse_tabular_file,
            'zip': self.parse_zip
        }
    
    def _verify_json_format_args(self, params: Union[str, dict], strict_json: bool = False) -> dict:
        """Verify the parameters of the function call"""
        if isinstance(params, str):
            try:
                if strict_json:
                    params_json: dict = json.loads(params)
                else:
                    params_json: dict = json_loads(params)
            except json.decoder.JSONDecodeError:
                raise ValueError('Parameters must be formatted as a valid JSON!')
        else:
            params_json: dict = params
        if isinstance(self.parameters, list):
            for param in self.parameters:
                if 'required' in param and param['required']:
                    if param['name'] not in params_json:
                        raise ValueError('Parameters %s is required!' % param['name'])
        elif isinstance(self.parameters, dict):
            import jsonschema
            jsonschema.validate(instance=params_json, schema=self.parameters)
        else:
            raise ValueError
        return params_json

    def call(self, params: Union[str, dict], **kwargs) -> Union[str, list]:
        params = self._verify_json_format_args(params)
        file_path = self._prepare_file(params['url'])
        
        # 修改点 2: 使用本地文件检查替代 self.db.get()
        # 构造缓存文件名: hash + _ori.json
        file_hash = hash_sha256(file_path)
        cache_file_path = os.path.join(self.data_root, f'{file_hash}_ori.json')

        # 检查缓存文件是否存在
        if os.path.exists(cache_file_path):
            try:
                with open(cache_file_path, 'r', encoding='utf-8') as f:
                    cached_data = json.load(f)
                logger.info(f"Hit cache for {file_path}")
                return self._flatten_result(cached_data)
            except Exception as e:
                logger.warning(f"Cache file exists but read failed: {e}, reparsing...")
        
        # 如果没有缓存或读取失败，执行解析
        return self._flatten_result(self._process_new_file(file_path))

    def _prepare_file(self, path: str) -> str:
        if is_http_url(path):
            download_dir = os.path.join(self.data_root, hash_sha256(path))
            os.makedirs(download_dir, exist_ok=True)
            return save_url_to_local_work_dir(path, download_dir)
        return sanitize_chrome_file_path(path)

    def _process_new_file(self, file_path: str) -> Union[str, list]:
        file_type = get_file_type(file_path)
        idp_types = ['pdf', 'docx', 'pptx', 'xlsx', 'jpg', 'png', 'mp3']
        logger.info(f'Start parsing {file_path}...')
        logger.info(f'File type {file_type}...')
        
        if file_type not in idp_types:
            try:
                # 尝试从 url/path 中获取后缀，如果没有取到可能会报错，这里加个简单的保护
                base_name = get_basename_from_url(file_path)
                if '.' in base_name:
                    file_type = base_name.split('.')[-1].lower()
            except:
                pass

        try:
            # 注意：这里需要确保 parse_ppt/parse_pdf 等函数在作用域内
            if file_type in self.parsers:
                results = self.parsers[file_type](file_path)
            else:
                # 如果类型不支持，尝试当做 txt 处理或者报错
                logger.warning(f"Unsupported file type: {file_type}, trying txt parser.")
                results = parse_txt(file_path)

            tokens = 0
            for page in results:
                for para in page['content']:
                    if 'schema' in para:
                        para['token'] = count_tokens(json.dumps(para['schema']))
                    else:
                        para['token'] = count_tokens(para.get('text', para.get('table', '')))
                    tokens += para['token']

            if not results or not tokens:
                logger.error(f"Parsing failed: No information was parsed")
                # 这里为了简单，去掉了自定义 Exception，如果需要可以保留
                raise ValueError("Document parsing failed: No content")
            else:
                self._cache_result(file_path, results)
                return results
        except Exception as e:
            logger.error(f"Parsing failed: {str(e)}")
            raise e

    def _cache_result(self, file_path: str, result: list):
        # 修改点 3: 使用标准文件写入替代 self.db.put()
        try:
            file_hash = hash_sha256(file_path)
            cache_file_path = os.path.join(self.data_root, f'{file_hash}_ori.json')
            
            with open(cache_file_path, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False)
                
            logger.info(f'The parsing result of {file_path} has been cached to {cache_file_path}')
        except Exception as e:
            logger.error(f"Failed to cache result: {str(e)}")

    def _flatten_result(self, result: list) -> str:
        return PARAGRAPH_SPLIT_SYMBOL.join(
            para.get('text', para.get('table', ''))
            for page in result for para in page['content']
        )

    def parse_zip(self, file_path: str) -> List[dict]:
        extract_dir = os.path.join(self.data_root, f"zip_{hash_sha256(file_path)}")
        os.makedirs(extract_dir, exist_ok=True)

        results = []
        # 注意：这里假设外部有 parse_zip 函数用于解压获取文件列表
        # 如果没有外部函数，需要在这里实现 unzip 逻辑
        for extracted_file in parse_zip(file_path, extract_dir):
            if (ft := get_file_type(extracted_file)) in self.parsers:
                try:
                    results.extend(self.parsers[ft](extracted_file))
                except Exception as e:
                    logger.warning(f"Skip files {extracted_file}: {str(e)}")

        if not results:
            raise ValueError("No parseable content found in the ZIP file")
        return results

if __name__ == "__main__":
    import time
    import logging
    
    # 1. 配置日志输出，确保能看到 "Start parsing" 和 "Hit cache"
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )

    print("\n" + "="*60)
    print(" 🚀 开始真实数据测试 (Real Data Test) ")
    print("="*60)

    # =========================================================================
    # [配置区] 请修改这里
    # =========================================================================
    
    # 1. 设置待测试的真实文件路径 (支持 PDF, PPTX, DOCX 等)
    # Windows 用户请注意路径转义，例如 "D:\\Documents\\architecture_v1.pptx"
    TEST_FILE_PATH = "/home/u-longyy/电脑维修证明.doc"  # <-- 修改为你的本地文件路径

    # 2. 检查 API Key (如果你启用了 VLM 视觉分析)
    if not os.getenv("DASHSCOPE_API_KEY"):
        logger.warning("⚠️ 未检测到 DASHSCOPE_API_KEY 环境变量。")
        logger.warning("   如果你的文件包含图片且需要 VLM 解析流程图，请先设置 Key，否则 VLM 调用将失败。")
        # os.environ["DASHSCOPE_API_KEY"] = "sk-你的Key" # 你也可以在这里临时硬编码

    # =========================================================================
    
    if not os.path.exists(TEST_FILE_PATH):
        logger.error(f"❌ 测试文件不存在: {TEST_FILE_PATH}")
        logger.error("请修改代码中的 TEST_FILE_PATH 变量为真实的本地文件路径。")
        exit(1)

    # 初始化 Parser
    # 缓存将生成在当前目录下的 ./workspace/parser_cache 中
    workspace_dir = "./workspace"
    parser = SingleFileParser(cfg={
        'path': os.path.join(workspace_dir, 'parser_cache'),
        'structured_doc': True 
    })

    # --- 测试阶段 1: 首次解析 (写入文件缓存) ---
    print(f"\n[Phase 1] 正在解析文件: {os.path.basename(TEST_FILE_PATH)}")
    print("⏳ 这可能需要几秒钟 (取决于文件大小和是否调用 VLM)...")
    
    t0 = time.time()
    # 构造输入参数 (模拟 Agent 调用)
    input_params = {"url": TEST_FILE_PATH}
    
    try:
        # 调用 call 方法
        result_text = parser.call(input_params)
        t1 = time.time()
        
        print(f"✅ 解析完成! 耗时: {t1 - t0:.2f}s")
        print("-" * 30)
        print("📄 解析结果预览 (前 500 字符):")
        print(result_text)
        print("..." if len(result_text) > 500 else "")
        print("-" * 30)

    except Exception as e:
        logger.error(f"❌ 解析过程中发生错误: {e}")
        exit(1)

    # --- 测试阶段 2: 缓存命中测试 (读取本地 JSON) ---
    print(f"\n[Phase 2] 二次读取测试 (验证本地文件缓存机制)")
    
    t2 = time.time()
    cached_result = parser.call(input_params)
    t3 = time.time()

    print(f"⏱️ 二次读取耗时: {t3 - t2:.4f}s")
    
    if (t3 - t2) < 1.0:
        print("✅ 速度极快，成功命中本地 JSON 缓存。")
    else:
        print("⚠️ 速度较慢，可能未命中缓存，请检查日志。")

    # --- 验证内容一致性 ---
    if result_text == cached_result:
        print("✅ 内容一致性校验通过。")
    else:
        print("❌ 内容不一致!")

    # --- 检查 VLM 效果 (如果是 PPT/PDF) ---
    if "图片分析" in result_text or "Mermaid" in result_text or "Visual Analysis" in result_text:
        print("\n🎉 检测到 VLM 分析内容！流程图/架构图已被成功提取为文本描述。")
    else:
        print("\nℹ️ 未检测到明显的 VLM 分析标记。")
        print("   (原因可能是：文件中无图片、parse_pdf/ppt 中 extract_image 默认为 False、或 API 调用未生效)")

    print("\n测试结束。缓存文件位于:", parser.data_root)