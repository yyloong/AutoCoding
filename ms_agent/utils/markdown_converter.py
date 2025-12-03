# flake8: noqa
# isort: skip_file
# yapf: disable
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Dict, List, Tuple

import markdown
from ms_agent.utils.logger import get_logger
from ms_agent.utils.utils import install_package, is_package_installed

logger = get_logger()


def _install_playwright() -> bool:
    """
    Automatically install Playwright and Chromium browser

    Returns:
        bool: Whether installation was successful
    """
    if is_package_installed('playwright'):
        return True

    try:
        logger.info('Automatically installing Playwright ...')

        # Install playwright package
        result1 = subprocess.run(
            [sys.executable, '-m', 'pip', 'install', 'playwright'],
            capture_output=True,
            text=True,
            timeout=300)

        if result1.returncode != 0:
            logger.error(
                f'Playwright package installation failed: {result1.stderr}')
            return False

        # Install Chromium browser
        result2 = subprocess.run(['playwright', 'install', 'chromium'],
                                 capture_output=True,
                                 text=True,
                                 timeout=600)

        if result2.returncode != 0:
            logger.error(
                f'Chromium browser installation failed: {result2.stderr}')
            # Try using python -m playwright install chromium
            logger.info('🔄 Trying alternative installation method...')
            result2_alt = subprocess.run(
                [sys.executable, '-m', 'playwright', 'install', 'chromium'],
                capture_output=True,
                text=True,
                timeout=600)

            if result2_alt.returncode != 0:
                logger.error(
                    f'Alternative method also failed: {result2_alt.stderr}')
                return False

        # Verify installation
        logger.info('Verifying Playwright installation...')
        try:
            from playwright.sync_api import sync_playwright
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                browser.close()
            logger.info('Playwright installation verification successful')
            return True
        except Exception as e:
            logger.error(f'Playwright verification failed: {e}')
            return False

    except subprocess.TimeoutExpired:
        logger.error(
            'Installation timeout, please install Playwright manually')
        return False
    except Exception as e:
        logger.error(f'Error during installation process: {e}')
        return False


class MarkdownConverter:
    """
    A class to convert Markdown files to HTML, PDF, DOCX, and PPT formats.

    Notes:
        This class requires the `markdown`, `beautifulsoup4`, `python-docx`, and `python-pptx` packages.
        It also uses Playwright for PDF conversion, which requires the Playwright package and Chromium browser.
        These packages will be installed automatically when calling the conversion methods.
    """

    @staticmethod
    def _extract_images_from_markdown(markdown_content: str) -> List[Tuple[str, str]]:
        """
        Extract image references from Markdown content

        Returns:
            List of (image_path, alt_text) tuples
        """
        # Match ![alt](path) format images
        image_pattern = r'!\[([^\]]*)\]\(([^)]+)\)'
        matches = re.findall(image_pattern, markdown_content)
        return [(match[1], match[0]) for match in matches]

    @staticmethod
    def _copy_and_process_images(
            images: List[Tuple[str, str]],
            output_dir: Path, folder_path: Path) -> Dict[str, str]:
        """
        Copy and process image files

        Returns:
            Dictionary mapping original paths to new paths
        """
        from PIL import Image

        image_mapping = {}
        output_images_dir = output_dir / 'images'
        output_images_dir.mkdir(exist_ok=True)

        resources_dir = folder_path / 'resources'

        for img_path, alt_text in images:
            # Handle relative paths
            if not os.path.isabs(img_path):
                # First try resources folder, then folder_path
                source_path = resources_dir / img_path
                if not source_path.exists():
                    source_path = folder_path / img_path
            else:
                source_path = Path(img_path)

            if source_path.exists():
                # Generate new filename
                new_filename = f'{len(image_mapping)}_{source_path.name}'
                new_path = output_images_dir / new_filename

                try:
                    # Copy and possibly compress image
                    if source_path.suffix.lower() in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                        with Image.open(source_path) as img:
                            # Compress if image is too large
                            if img.width > 1200 or img.height > 1200:
                                img.thumbnail((1200, 1200),
                                              Image.Resampling.LANCZOS)
                            img.save(new_path, optimize=True, quality=85)
                    else:
                        shutil.copy2(source_path, new_path)

                    image_mapping[img_path] = str(
                        new_path.relative_to(output_dir))
                except (IOError, OSError) as e:
                    logger.warning(f'Error processing image {img_path}: {e}')
                    image_mapping[img_path] = img_path
            else:
                logger.warning(f'Image file not found: {source_path}')
                image_mapping[img_path] = img_path

        return image_mapping

    @staticmethod
    def _update_markdown_image_paths(markdown_content: str,
                                     image_mapping: Dict[str, str]) -> str:
        """
        Update image paths in Markdown content
        """

        def replace_image_path(match):
            alt_text = match.group(1)
            original_path = match.group(2)
            new_path = image_mapping.get(original_path, original_path)
            return f'![{alt_text}]({new_path})'

        image_pattern = r'!\[([^\]]*)\]\(([^)]+)\)'
        return re.sub(image_pattern, replace_image_path, markdown_content)

    @staticmethod
    def _add_formatted_text(paragraph, text):
        """
        Add formatted text to a paragraph, handling basic Markdown formatting
        """
        # Handle bold (**text** or __text__)
        bold_pattern = r'\*\*(.*?)\*\*|__(.*?)__'
        # Handle italic (*text* or _text_)
        italic_pattern = r'(?<!\*)\*([^*]+)\*(?!\*)|(?<!_)_([^_]+)_(?!_)'
        # Handle inline code (`code`)
        code_pattern = r'`([^`]+)`'

        # Split text by formatting patterns
        parts = []
        last_end = 0

        # Find all formatting matches
        for match in re.finditer(
                f'({bold_pattern})|({italic_pattern})|({code_pattern})', text):
            # Add text before match
            if match.start() > last_end:
                parts.append(('text', text[last_end:match.start()]))

            # Determine match type and content
            if match.group(2) or match.group(3):  # Bold
                content = match.group(2) or match.group(3)
                parts.append(('bold', content))
            elif match.group(5) or match.group(6):  # Italic
                content = match.group(5) or match.group(6)
                parts.append(('italic', content))
            elif match.group(8):  # Code
                content = match.group(8)
                parts.append(('code', content))

            last_end = match.end()

        # Add remaining text
        if last_end < len(text):
            parts.append(('text', text[last_end:]))

        # If no formatting found, add as plain text
        if not parts:
            parts = [('text', text)]

        # Add formatted runs to paragraph
        for part_type, content in parts:
            run = paragraph.add_run(content)
            if part_type == 'bold':
                run.bold = True
            elif part_type == 'italic':
                run.italic = True
            elif part_type == 'code':
                run.font.name = 'Courier New'

    @staticmethod
    def _process_html_inline_elements(paragraph, element):
        """
        Process HTML inline elements and add formatted runs to paragraph
        """
        # Handle mixed content (text and inline elements)
        for content in element.contents:
            if hasattr(content, 'name'):  # It's a tag
                if content.name == 'strong' or content.name == 'b':
                    run = paragraph.add_run(content.get_text())
                    run.bold = True
                elif content.name == 'em' or content.name == 'i':
                    run = paragraph.add_run(content.get_text())
                    run.italic = True
                elif content.name == 'code':
                    run = paragraph.add_run(content.get_text())
                    run.font.name = 'Courier New'
                elif content.name == 'a':
                    # For links, just add the text (DOCX hyperlinks are complex)
                    run = paragraph.add_run(content.get_text())
                else:
                    # For other tags, just add the text content
                    run = paragraph.add_run(content.get_text())
            else:  # It's text content
                text_content = str(content).strip()
                if text_content:
                    paragraph.add_run(text_content)

    @staticmethod
    def _process_html_table(doc, table_element):
        """
        Process HTML table element and create DOCX table
        """
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        try:
            # Extract table data from HTML
            rows = []
            for tr in table_element.find_all('tr'):
                row_data = []
                for cell in tr.find_all(['td', 'th']):
                    cell_text = cell.get_text().strip()
                    row_data.append(cell_text)
                if row_data:  # Only add non-empty rows
                    rows.append(row_data)

            if not rows:
                return

            # Determine table dimensions
            max_cols = max(len(row) for row in rows)
            num_rows = len(rows)

            if max_cols > 0 and num_rows > 0:
                # Create DOCX table
                table = doc.add_table(rows=num_rows, cols=max_cols)
                table.style = 'Table Grid'
                table.autofit = False

                # Fill table data
                for row_idx, row_data in enumerate(rows):
                    table_row = table.rows[row_idx]

                    for col_idx in range(max_cols):
                        cell = table_row.cells[col_idx]

                        # Get cell data or empty string if not enough columns
                        cell_data = row_data[col_idx] if col_idx < len(row_data) else ''

                        # Clear existing content and add new
                        cell.text = ''
                        paragraph = cell.paragraphs[0]
                        paragraph.add_run(cell_data)

                        # Make header row bold and centered (first row or if original was th)
                        original_tr = table_element.find_all('tr')[row_idx]
                        if row_idx == 0 or original_tr.find('th'):
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            for run in paragraph.runs:
                                run.bold = True

                # Add spacing after table
                doc.add_paragraph()

        except Exception as e:
            logger.warning(f'Error processing HTML table: {e}')
            # Fallback: add table as text
            table_text = table_element.get_text().strip()
            if table_text:
                doc.add_paragraph(table_text)

    @staticmethod
    def _process_markdown_table_fallback(doc, table_lines):
        """
        Fallback method to process markdown table syntax directly
        """
        try:
            # Parse table structure from markdown lines
            table_data = []

            for line in table_lines:
                # Clean up the line and split by |
                cells = [cell.strip() for cell in line.split('|')]

                # Remove empty cells at start/end (from leading/trailing |)
                while cells and not cells[0]:
                    cells.pop(0)
                while cells and not cells[-1]:
                    cells.pop()

                # Skip empty rows and separator lines
                if not cells:
                    continue

                # Check if this is a separator line (contains only -, :, |, and spaces)
                is_separator = all(
                    re.match(r'^[-:\s|]*$', cell) and ('-' in cell or ':' in cell)
                    for cell in cells if cell.strip()
                )

                if is_separator:
                    continue

                # This is a data row
                if cells:
                    table_data.append(cells)

            if table_data:
                # Create DOCX table
                max_cols = max(len(row) for row in table_data)
                num_rows = len(table_data)

                if max_cols > 0 and num_rows > 0:
                    table = doc.add_table(rows=num_rows, cols=max_cols)
                    table.style = 'Table Grid'

                    # Fill table data
                    for row_idx, row_data in enumerate(table_data):
                        table_row = table.rows[row_idx]

                        for col_idx in range(max_cols):
                            cell = table_row.cells[col_idx]
                            cell_data = row_data[col_idx] if col_idx < len(row_data) else ''
                            cell.text = cell_data

                            # Make first row bold (header)
                            if row_idx == 0:
                                for paragraph in cell.paragraphs:
                                    for run in paragraph.runs:
                                        run.bold = True

                    # Add spacing after table
                    doc.add_paragraph()

        except Exception as e:
            logger.warning(f'Error processing markdown table fallback: {e}')
            # Add table as text if processing fails
            table_text = '\n'.join(table_lines)
            doc.add_paragraph(table_text)

    @staticmethod
    def _process_docx_image(doc, img_src, img_alt, resources_dir, markdown_folder):
        """
        Process image and add to DOCX document
        """
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches as DocxInches

        try:
            # Handle image path
            if not os.path.isabs(img_src):
                actual_img_path = resources_dir / img_src
                if not actual_img_path.exists():
                    actual_img_path = markdown_folder / img_src
            else:
                actual_img_path = Path(img_src)

            if actual_img_path.exists():
                try:
                    # Add image caption before the image if it exists
                    if img_alt:
                        caption_paragraph = doc.add_paragraph()
                        caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        caption_run = caption_paragraph.add_run(img_alt)
                        caption_run.italic = True
                        caption_run.font.size = DocxInches(0.12)

                    # Add the image
                    paragraph = doc.add_paragraph()
                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                    run.add_picture(str(actual_img_path), width=DocxInches(6))
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

                except Exception as e:
                    doc.add_paragraph(f'[Image loading failed: {img_src}]')
                    logger.warning(f'Failed to add image: {e}')
            else:
                logger.warning(f'Image file not found: {actual_img_path}')
                doc.add_paragraph(f'[Image not found: {img_src}]')

        except Exception as e:
            logger.warning(f'Error processing image {img_src}: {e}')
            doc.add_paragraph(f'[Image processing failed: {img_src}]')

    @staticmethod
    def markdown_to_html(markdown_path: str, html_path: str, lang: str = 'zh-CN') -> str:
        """
        Convert Markdown files to HTML

        Args:
            markdown_path: Path to folder containing Markdown files and resources, or path to a single .md file
            html_path: Path to output HTML folder or file
            lang: Language code for the HTML document (default: 'zh-CN'), e.g. 'zh-CN', 'en-US'

        Returns:
            Absolute path to generated HTML folder or file
        """
        install_package(package_name='beautifulsoup4', import_name='bs4')
        from bs4 import BeautifulSoup

        markdown_input = Path(markdown_path)

        # Determine if input is a file or folder
        if markdown_input.is_file() and markdown_input.suffix == '.md':
            # Single file mode
            markdown_files = [markdown_input]
            markdown_folder = markdown_input.parent

            # If html_path ends with .html, treat as single file output
            html_output = Path(html_path)
            if html_output.suffix == '.html':
                html_folder = html_output.parent
                html_folder.mkdir(parents=True, exist_ok=True)
                single_file_output = html_output
            else:
                html_folder = html_output
                html_folder.mkdir(parents=True, exist_ok=True)
                single_file_output = None
        else:
            # Folder mode
            markdown_folder = markdown_input
            html_folder = Path(html_path)
            html_folder.mkdir(parents=True, exist_ok=True)
            single_file_output = None

            # Find all .md files in the markdown folder
            markdown_files = list(markdown_folder.glob('*.md'))
            if not markdown_files:
                raise FileNotFoundError(
                    f'No .md files found in {markdown_folder}')

        logger.info(f'Found {len(markdown_files)} Markdown files to convert to html:')
        for md_file in markdown_files:
            logger.info(f'  - {md_file.name}')

        # Process each Markdown file
        converted_files = []
        for markdown_file in markdown_files:
            try:
                # Read Markdown file
                with open(markdown_file, 'r', encoding='utf-8') as f:
                    markdown_content = f.read()

                # Extract and process images for this file
                images = MarkdownConverter._extract_images_from_markdown(
                    markdown_content)
                image_mapping = MarkdownConverter._copy_and_process_images(
                    images, html_folder, markdown_folder)
                updated_markdown = MarkdownConverter._update_markdown_image_paths(
                    markdown_content, image_mapping)

                # Convert to HTML
                md = markdown.Markdown(
                    extensions=['extra', 'codehilite', 'toc'])
                html_content = md.convert(updated_markdown)

                # Get the title from the first H1 heading or use filename
                title = markdown_file.stem
                soup = BeautifulSoup(html_content, 'html.parser')
                h1_tag = soup.find('h1')
                if h1_tag:
                    title = h1_tag.get_text()

                # Create complete HTML document
                full_html = f"""
<!DOCTYPE html>
<html lang="{lang}">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            line-height: 1.6;
            color: #333;
            max-width: 800px;
            margin: 0 auto;
            padding: 20px;
        }}
        img {{
            max-width: 100%;
            height: auto;
            border-radius: 8px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        }}
        pre {{
            background: #f4f4f4;
            padding: 15px;
            border-radius: 5px;
            overflow-x: auto;
        }}
        blockquote {{
            border-left: 4px solid #ddd;
            margin: 0;
            padding-left: 20px;
            color: #666;
        }}
        table {{
            border-collapse: collapse;
            width: 100%;
        }}
        th, td {{
            border: 1px solid #ddd;
            padding: 12px;
            text-align: left;
        }}
        th {{
            background-color: #f2f2f2;
        }}
    </style>
</head>
<body>
    {html_content}
</body>
</html>
                """

                # Generate output HTML file path
                if single_file_output:
                    output_html_file = single_file_output
                else:
                    output_html_file = html_folder / f'{markdown_file.stem}.html'

                # Write HTML file
                with open(output_html_file, 'w', encoding='utf-8') as f:
                    f.write(full_html)

                converted_files.append(output_html_file)
                logger.info(f'✓ HTML file generated: {output_html_file}')

            except Exception as e:
                logger.info(f'✗ Failed to convert {markdown_file.name}: {e}')

        logger.info(
            f'\nConversion completed! Generated {len(converted_files)} HTML files.'
        )
        if single_file_output:
            return str(single_file_output.absolute())
        else:
            return str(html_folder.absolute())

    @staticmethod
    def markdown_to_pdf(markdown_path: str, pdf_path: str) -> str:
        """
        Convert Markdown files to PDF

        Args:
            markdown_path: Path to folder containing Markdown files and resources, or path to a single .md file
            pdf_path: Path to output PDF folder or file

        Returns:
            Absolute path to generated PDF folder or file
        """
        _install_playwright()

        try:
            from playwright.sync_api import sync_playwright
        except ImportError:
            logger.error(
                'Playwright is not installed. Please install it using:\n'
                '   pip install playwright\n'
                '   playwright install chromium')
            return ''

        import tempfile

        markdown_input = Path(markdown_path)

        # Determine if input is a file or folder
        if markdown_input.is_file() and markdown_input.suffix == '.md':
            # Single file mode
            markdown_files = [markdown_input]
            markdown_folder = markdown_input.parent

            # If pdf_path ends with .pdf, treat as single file output
            pdf_output = Path(pdf_path)
            if pdf_output.suffix == '.pdf':
                pdf_folder = pdf_output.parent
                pdf_folder.mkdir(parents=True, exist_ok=True)
                single_file_output = pdf_output
            else:
                pdf_folder = pdf_output
                pdf_folder.mkdir(parents=True, exist_ok=True)
                single_file_output = None
        else:
            # Folder mode
            markdown_folder = markdown_input
            pdf_folder = Path(pdf_path)
            pdf_folder.mkdir(parents=True, exist_ok=True)
            single_file_output = None

            # Find all .md files in the markdown folder
            markdown_files = list(markdown_folder.glob('*.md'))
            if not markdown_files:
                raise FileNotFoundError(
                    f'No .md files found in {markdown_folder}')

        logger.info(f'Found {len(markdown_files)} Markdown files to convert to pdf:')
        for md_file in markdown_files:
            logger.info(f'  - {md_file.name}')

        # Process each Markdown file
        # TODO: can be multi-threaded for faster processing
        converted_files = []
        for markdown_file in markdown_files:
            try:
                # First convert to HTML
                with tempfile.TemporaryDirectory() as temp_dir_str:
                    temp_dir = Path(temp_dir_str)

                    # Convert this specific markdown file to HTML first
                    MarkdownConverter.markdown_to_html(str(markdown_file), str(temp_dir))
                    temp_html = temp_dir / f'{markdown_file.stem}.html'

                    # Generate output PDF file path
                    if single_file_output:
                        output_pdf_file = single_file_output
                    else:
                        output_pdf_file = pdf_folder / f'{markdown_file.stem}.pdf'

                    try:
                        logger.info(
                            f'Converting {markdown_file.name} to PDF using Playwright...'
                        )

                        with sync_playwright() as p:
                            # Launch browser (prefer Chromium for best PDF support)
                            browser = p.chromium.launch(headless=True)
                            page = browser.new_page()

                            # Enhanced PDF options for high quality output
                            pdf_options = {
                                'path': str(output_pdf_file),
                                'format': 'A4',
                                'margin': {
                                    'top': '2cm',
                                    'right': '2cm',
                                    'bottom': '2cm',
                                    'left': '2cm'
                                },
                                'print_background':
                                    True,
                                'prefer_css_page_size':
                                    True,
                                'display_header_footer':
                                    True,
                                'header_template':
                                    '<div style="font-size:10px; text-align:center; width:100%;"></div>',
                                'footer_template':
                                    '<div style="font-size:10px; text-align:center; width:100%; margin-top:10px;"><span class="pageNumber"></span> / <span class="totalPages"></span></div>'
                            }

                            # Navigate to HTML file and wait for full load
                            page.goto(
                                f'file://{temp_html.absolute()}',
                                wait_until='networkidle')

                            # Wait a bit more for any dynamic content
                            page.wait_for_timeout(1000)

                            # Generate PDF
                            page.pdf(**pdf_options)

                            browser.close()

                        logger.info(
                            f'PDF file generated using Playwright: {output_pdf_file}'
                        )
                        converted_files.append(output_pdf_file)

                    except Exception as e:
                        logger.error(
                            f'Playwright conversion failed for {markdown_file.name}: {e}'
                        )

            except Exception as e:
                logger.error(f'✗ Failed to convert {markdown_file.name}: {e}')

        logger.info(
            f'\nConversion completed! Generated {len(converted_files)} PDF files.'
        )
        if single_file_output:
            return str(single_file_output.absolute())
        else:
            return str(pdf_folder.absolute())

    @staticmethod
    def markdown_to_docx(markdown_path: str, docx_path: str) -> str:
        """
        Convert Markdown files to DOCX

        Args:
            markdown_path: Path to folder containing Markdown files and resources, or path to a single .md file
            docx_path: Path to output DOCX folder or file

        Returns:
            Absolute path to generated DOCX folder or file
        """
        # Check and install the `python-docx` and `beautifulsoup4` automatically
        install_package(package_name='python-docx', import_name='docx')
        install_package(package_name='beautifulsoup4', import_name='bs4')

        from bs4 import BeautifulSoup
        from docx import Document

        markdown_input = Path(markdown_path)

        # Determine if input is a file or folder
        if markdown_input.is_file() and markdown_input.suffix == '.md':
            # Single file mode
            markdown_files = [markdown_input]
            markdown_folder = markdown_input.parent

            # If docx_path ends with .docx, treat as single file output
            docx_output = Path(docx_path)
            if docx_output.suffix == '.docx':
                docx_folder = docx_output.parent
                docx_folder.mkdir(parents=True, exist_ok=True)
                single_file_output = docx_output
            else:
                docx_folder = docx_output
                docx_folder.mkdir(parents=True, exist_ok=True)
                single_file_output = None
        else:
            # Folder mode
            markdown_folder = markdown_input
            docx_folder = Path(docx_path)
            docx_folder.mkdir(parents=True, exist_ok=True)
            single_file_output = None

            # Find all .md files in the markdown folder
            markdown_files = list(markdown_folder.glob('*.md'))
            if not markdown_files:
                raise FileNotFoundError(
                    f'No .md files found in {markdown_folder}')

        logger.info(f'Found {len(markdown_files)} Markdown files to convert to docx:')
        for md_file in markdown_files:
            logger.info(f'  - {md_file.name}')

        resources_dir = markdown_folder / 'resources'

        # Process each Markdown file
        converted_files = []
        for markdown_file in markdown_files:
            try:
                # Read Markdown file
                with open(markdown_file, 'r', encoding='utf-8') as f:
                    markdown_content = f.read()

                # Create Word document (images will be embedded directly)
                doc = Document()

                # Use proper Markdown parsing with enhanced extensions
                md = markdown.Markdown(extensions=[
                    'extra',  # Includes tables, fenced_code, footnotes, attr_list, def_list, abbr
                    'codehilite',  # Code highlighting
                    'toc',  # Table of contents
                    'tables',  # Explicit table support
                    'fenced_code',  # Fenced code blocks
                    'attr_list'  # Attribute lists
                ])
                html_content = md.convert(markdown_content)

                # Parse the HTML using BeautifulSoup for structured processing
                soup = BeautifulSoup(html_content, 'html.parser')

                # First, handle images that might not be converted to <img> tags by markdown
                # Process raw markdown image syntax that might have been missed
                remaining_content = str(soup)
                img_pattern = r'!\[([^\]]*)\]\(([^)]+)\)'
                img_matches = re.findall(img_pattern, remaining_content)

                # Add missed images to the soup
                for alt_text, img_src in img_matches:
                    # Create img element and add to soup
                    img_tag = soup.new_tag('img', src=img_src, alt=alt_text)
                    # Find a good place to insert it (after last processed element or at end)
                    if soup.body:
                        soup.body.append(img_tag)
                    else:
                        soup.append(img_tag)

                # Process each HTML element and convert to DOCX
                for element in soup.find_all(
                        ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'blockquote', 'pre', 'table', 'img']):
                    try:
                        if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                            # Headings
                            level = int(element.name[1])
                            if level <= 3:  # DOCX supports up to 9 levels, but we'll use 1-3 for clarity
                                doc.add_heading(element.get_text().strip(), level=level)
                            else:
                                # For h4-h6, use bold paragraph instead
                                paragraph = doc.add_paragraph()
                                run = paragraph.add_run(element.get_text().strip())
                                run.bold = True

                        elif element.name == 'p':
                            # Paragraphs - handle inline formatting
                            if element.get_text().strip():  # Skip empty paragraphs
                                paragraph = doc.add_paragraph()
                                MarkdownConverter._process_html_inline_elements(paragraph, element)

                        elif element.name == 'ul':
                            # Unordered lists
                            for li in element.find_all('li', recursive=False):
                                paragraph = doc.add_paragraph(style='List Bullet')
                                MarkdownConverter._process_html_inline_elements(paragraph, li)

                        elif element.name == 'ol':
                            # Ordered lists
                            for li in element.find_all('li', recursive=False):
                                paragraph = doc.add_paragraph(style='List Number')
                                MarkdownConverter._process_html_inline_elements(paragraph, li)

                        elif element.name == 'blockquote':
                            # Blockquotes
                            for p in element.find_all('p'):
                                paragraph = doc.add_paragraph(style='Quote')
                                MarkdownConverter._process_html_inline_elements(paragraph, p)

                        elif element.name == 'pre':
                            # Code blocks
                            code_element = element.find('code')
                            if code_element:
                                code_text = code_element.get_text()
                            else:
                                code_text = element.get_text()
                            paragraph = doc.add_paragraph(code_text)
                            # Apply code formatting
                            for run in paragraph.runs:
                                run.font.name = 'Courier New'

                        elif element.name == 'table':
                            # Tables - enhanced processing
                            MarkdownConverter._process_html_table(doc, element)

                        elif element.name == 'img':
                            # Images
                            img_src = element.get('src', '')
                            img_alt = element.get('alt', '')
                            if img_src:  # Only process if src exists
                                MarkdownConverter._process_docx_image(doc, img_src, img_alt, resources_dir,
                                                                      markdown_folder)


                    except Exception as e:
                        logger.warning(f'Error processing element {element.name}: {e}')
                        # Fallback: add as plain text
                        if element.get_text().strip():
                            doc.add_paragraph(element.get_text().strip())

                # Additional fallback processing for missed elements
                # Handle any remaining markdown table syntax that wasn't converted
                if '|' in markdown_content:
                    lines = markdown_content.split('\n')
                    table_lines = []
                    in_table = False

                    for line in lines:
                        line = line.strip()
                        if '|' in line and line:
                            if not in_table:
                                in_table = True
                                table_lines = []
                            table_lines.append(line)
                        else:
                            if in_table and table_lines:
                                # Process collected table
                                MarkdownConverter._process_markdown_table_fallback(doc, table_lines)
                                table_lines = []
                                in_table = False

                    # Process last table if exists
                    if in_table and table_lines:
                        MarkdownConverter._process_markdown_table_fallback(doc, table_lines)

                # Handle any remaining markdown image syntax that wasn't converted
                remaining_img_matches = re.findall(r'!\[([^\]]*)\]\(([^)]+)\)', markdown_content)
                for alt_text, img_src in remaining_img_matches:
                    # Check if this image was already processed
                    already_processed = False
                    for element in soup.find_all('img'):
                        if element.get('src') == img_src:
                            already_processed = True
                            break

                    if not already_processed:
                        MarkdownConverter._process_docx_image(doc, img_src, alt_text, resources_dir, markdown_folder)

                # Generate output DOCX file path
                if single_file_output:
                    output_docx_file = single_file_output
                else:
                    output_docx_file = docx_folder / f'{markdown_file.stem}.docx'

                # Save document
                doc.save(output_docx_file)
                converted_files.append(output_docx_file)
                logger.info(f'✓ DOCX file generated: {output_docx_file}')

            except Exception as e:
                logger.info(f'✗ Failed to convert {markdown_file.name}: {e}')

        logger.info(
            f'\nConversion completed! Generated {len(converted_files)} DOCX files.'
        )
        if single_file_output:
            return str(single_file_output.absolute())
        else:
            return str(docx_folder.absolute())

    @staticmethod
    def markdown_to_pptx(markdown_path: str, ppt_path: str) -> str:
        """
        Convert Markdown files to PPT

        Args:
            markdown_path: Path to folder containing Markdown files and resources, or path to a single .md file
            ppt_path: Path to output PPT folder or file

        Returns:
            Absolute path to generated PPT folder or file
        """
        # Check and install the `python-pptx` automatically
        install_package(package_name='python-pptx', import_name='pptx')

        from pptx import Presentation
        from pptx.util import Inches as PptxInches
        from pptx.enum.text import PP_ALIGN

        markdown_input = Path(markdown_path)

        # Determine if input is a file or folder
        if markdown_input.is_file() and markdown_input.suffix == '.md':
            # Single file mode
            markdown_files = [markdown_input]
            markdown_folder = markdown_input.parent

            # If ppt_path ends with .pptx, treat as single file output
            ppt_output = Path(ppt_path)
            if ppt_output.suffix == '.pptx':
                ppt_folder = ppt_output.parent
                ppt_folder.mkdir(parents=True, exist_ok=True)
                single_file_output = ppt_output
            else:
                ppt_folder = ppt_output
                ppt_folder.mkdir(parents=True, exist_ok=True)
                single_file_output = None
        else:
            # Folder mode
            markdown_folder = markdown_input
            ppt_folder = Path(ppt_path)
            ppt_folder.mkdir(parents=True, exist_ok=True)
            single_file_output = None

            # Find all .md files in the markdown folder
            markdown_files = list(markdown_folder.glob('*.md'))
            if not markdown_files:
                raise FileNotFoundError(
                    f'No .md files found in {markdown_folder}')

        logger.info(f'Found {len(markdown_files)} Markdown files to convert to ppt:')
        for md_file in markdown_files:
            logger.info(f'  - {md_file.name}')

        resources_dir = markdown_folder / 'resources'

        # Process each Markdown file
        converted_files = []
        for markdown_file in markdown_files:
            try:
                # Read Markdown file
                with open(markdown_file, 'r', encoding='utf-8') as f:
                    markdown_content = f.read()

                # Create PPT
                prs = Presentation()

                # Parse content into slides
                slides_content = []
                current_slide = {'title': '', 'content': []}

                lines = markdown_content.split('\n')

                # Parse content into slides with table support
                i = 0
                while i < len(lines):
                    line = lines[i].strip()

                    if line.startswith('# '):
                        # New slide title - only add previous slide if it has meaningful content
                        if current_slide['title'] and current_slide['content']:
                            slides_content.append(current_slide)
                        elif current_slide['title'] and not current_slide['content']:
                            # Previous slide has title but no content, skip it
                            logger.info(
                                f"Skipping empty slide: {current_slide['title']}"
                            )
                        current_slide = {'title': line[2:], 'content': []}
                    elif line.startswith('## '):
                        # Level 2 heading as content
                        current_slide['content'].append(('heading', line[3:]))
                    elif line.startswith('!['):
                        # Image
                        img_match = re.match(r'!\[([^\]]*)\]\(([^)]+)\)', line)
                        if img_match:
                            alt_text, img_path = img_match.groups()
                            current_slide['content'].append(
                                ('image', img_path, alt_text))
                    elif '|' in line:
                        # Table detection - collect all consecutive table lines
                        table_lines = []
                        table_start = i

                        # Collect all consecutive table lines
                        while i < len(lines):
                            current_line = lines[i].strip()
                            if '|' in current_line:
                                table_lines.append(current_line)
                                i += 1
                            else:
                                break

                        # Process the table if we found lines
                        if table_lines:
                            # Parse table structure
                            table_data = []

                            for table_line in table_lines:
                                # Clean up the line and split by |
                                cells = [
                                    cell.strip()
                                    for cell in table_line.split('|')
                                ]

                                # Remove empty cells at start/end (from leading/trailing |)
                                while cells and not cells[0]:
                                    cells.pop(0)
                                while cells and not cells[-1]:
                                    cells.pop()

                                # Skip empty rows
                                if not cells:
                                    continue

                                # Check if this is a separator line (contains only -, :, |, and spaces)
                                is_separator = all(
                                    re.match(r'^[-:\s|]*$', cell) and ('-' in cell or ':' in cell) for cell in cells if cell.strip())

                                if is_separator:
                                    # This is a separator line, skip it
                                    continue

                                # This is a data row
                                if cells:
                                    table_data.append(cells)

                            # Add table to slide content if we have data
                            if table_data:
                                current_slide['content'].append(
                                    ('table', table_data))

                        # Continue from where we left off (i is already incremented)
                        continue
                    elif line.startswith('- ') or line.startswith('* '):
                        # List item
                        current_slide['content'].append(('bullet', line[2:]))
                    elif line and not line.startswith('#'):
                        # Regular text
                        current_slide['content'].append(('text', line))

                    i += 1

                # Add the last slide - only if it has meaningful content
                if current_slide['title'] and current_slide['content']:
                    slides_content.append(current_slide)
                elif current_slide['title'] and not current_slide['content']:
                    # Last slide has title but no content, skip it
                    logger.info(
                        f"Skipping empty slide: {current_slide['title']}")

                # Create slides with intelligent pagination
                for slide_data in slides_content:
                    # Skip slides without title or content
                    if not slide_data['title'] and not slide_data['content']:
                        continue

                    # Skip slides with only title but no content
                    if slide_data['title'] and not slide_data['content']:
                        logger.info(
                            f"Skipping slide with title only: {slide_data['title']}"
                        )
                        continue

                    # Get slide dimensions for calculations
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height

                    # Calculate adaptive font sizes based on slide dimensions
                    base_font_size = max(PptxInches(0.15), slide_width
                                         / 50)  # Adaptive base font
                    heading_font_size = base_font_size * 1.3  # 30% larger for headings
                    bullet_font_size = base_font_size * 0.9  # 10% smaller for bullets

                    # Calculate available content area (excluding title)
                    title_height = PptxInches(
                        1.5) if slide_data['title'] else PptxInches(0.5)
                    available_height = slide_height - title_height - PptxInches(
                        0.5)  # Bottom margin

                    # Estimate lines per slide based on font size and available height
                    line_height = base_font_size * 1.4  # Line spacing factor
                    max_lines_per_slide = int(
                        available_height / line_height) - 2  # Safety margin
                    max_lines_per_slide = max(
                        5, max_lines_per_slide)  # Minimum 5 lines per slide

                    # Group content into pages with advanced sparse page filtering
                    content_pages = []
                    current_page = []
                    current_line_count = 0
                    current_content_weight = 0  # Track content density

                    # Enhanced pagination parameters
                    min_content_per_page = 4  # Increased minimum content items per page
                    min_content_weight = 8  # Minimum content weight (text length + items)
                    optimal_lines_per_slide = max(
                        10, max_lines_per_slide // 2)  # More content per slide

                    # Process content and estimate line usage with content weight
                    for content_type, *content_data in slide_data['content']:
                        estimated_lines = 1  # Default for most content types
                        content_weight = 1  # Default content weight

                        if content_type in ['text', 'bullet', 'heading']:
                            text_content = content_data[0]
                            # More conservative line estimation
                            chars_per_line = int(
                                slide_width /
                                (base_font_size * 0.4))  # Even more generous
                            estimated_lines = max(1, len(text_content) // chars_per_line + 1)  # noqa

                            # Calculate content weight based on text length and type
                            content_weight = len(text_content) // 20 + 1  # noqa
                            if content_type == 'heading':
                                content_weight += 2  # Headings have more visual weight
                            elif content_type == 'bullet':
                                content_weight += 1  # Bullets have moderate weight

                            # Reduce extra spacing for headings
                            if content_type == 'heading':
                                estimated_lines = max(1, estimated_lines)

                        elif content_type == 'table':
                            # Tables need separate slides, don't include in text pagination
                            if current_page:
                                content_pages.append(current_page)
                                current_page = []
                                current_line_count = 0
                                current_content_weight = 0
                            # Add table as separate page
                            content_pages.append([(content_type, *content_data)
                                                  ])
                            continue

                        elif content_type == 'image':
                            # Images need separate slides, don't include in text pagination
                            if current_page:
                                content_pages.append(current_page)
                                current_page = []
                                current_line_count = 0
                                current_content_weight = 0
                            # Add image as separate page
                            content_pages.append([(content_type, *content_data)
                                                  ])
                            continue

                        # Enhanced page breaking logic with content weight consideration
                        should_break_page = False

                        # Break page conditions (all must be met for normal break):
                        # 1. Sufficient content items and weight
                        # 2. Would exceed optimal line limit
                        # 3. Current page has substantial content
                        if (current_line_count + estimated_lines >
                                optimal_lines_per_slide
                                and len(current_page) >= min_content_per_page
                                and
                                current_content_weight >= min_content_weight
                                and
                                current_line_count >= optimal_lines_per_slide
                                // 3):  # More lenient threshold
                            should_break_page = True

                        # Force break only if we're significantly over the limit
                        elif (current_line_count + estimated_lines >
                              max_lines_per_slide * 1.2 or
                              current_content_weight > min_content_weight * 2):
                            should_break_page = True

                        if should_break_page and current_page:
                            # Start new page
                            content_pages.append(current_page)
                            current_page = [(content_type, *content_data)]
                            current_line_count = estimated_lines
                            current_content_weight = content_weight
                        else:
                            # Add to current page
                            current_page.append((content_type, *content_data))
                            current_line_count += estimated_lines
                            current_content_weight += content_weight

                    # Add remaining content as last page
                    if current_page:
                        content_pages.append(current_page)

                    # Advanced post-processing: aggressive sparse page filtering and merging
                    filtered_pages = []
                    i = 0
                    while i < len(content_pages):
                        current_page = content_pages[i]

                        # Skip special content pages (tables/images) - they're always valid
                        if (len(current_page) == 1
                                and current_page[0][0] in ['table', 'image']):
                            filtered_pages.append(current_page)
                            i += 1
                            continue

                        # Calculate page content metrics
                        page_content_weight = 0
                        page_text_length = 0
                        for content_type, *content_data in current_page:
                            if content_type in ['text', 'bullet', 'heading']:
                                text_content = content_data[0]
                                page_text_length += len(text_content)
                                page_content_weight += len(
                                    text_content) // 20 + 1
                                if content_type == 'heading':
                                    page_content_weight += 2
                                elif content_type == 'bullet':
                                    page_content_weight += 1

                        # Determine if page is too sparse
                        is_sparse = (
                                len(current_page) < min_content_per_page
                                or page_content_weight < min_content_weight
                                or page_text_length < 100
                        )  # Less than 100 characters total

                        # Try to merge sparse pages with adjacent pages
                        if is_sparse and i + 1 < len(content_pages):
                            next_page = content_pages[i + 1]

                            # Don't merge with special content pages
                            if (len(next_page) == 1
                                    and next_page[0][0] in ['table', 'image']):
                                # Can't merge with special content, check if current page is acceptable
                                if len(current_page
                                       ) >= 2 or page_text_length >= 50:
                                    filtered_pages.append(current_page)
                                else:
                                    logger.info(
                                        f'Filtering out sparse page with '
                                        f'{len(current_page)} items, {page_text_length} chars'
                                    )
                                i += 1
                                continue

                            # Calculate combined metrics
                            next_page_weight = 0
                            next_page_text_length = 0
                            for content_type, *content_data in next_page:
                                if content_type in ['text', 'bullet', 'heading']:
                                    text_content = content_data[0]
                                    next_page_text_length += len(text_content)
                                    next_page_weight += len(
                                        text_content) // 20 + 1

                            combined_weight = page_content_weight + next_page_weight
                            combined_length = page_text_length + next_page_text_length
                            combined_items = len(current_page) + len(next_page)

                            # Merge if combined page would be reasonable
                            if (combined_items <= min_content_per_page * 2 and
                                    combined_weight <= min_content_weight * 2
                                    and combined_length <=
                                    800):  # Not too long when combined

                                merged_page = current_page + next_page
                                filtered_pages.append(merged_page)
                                logger.info(
                                    f'Merged sparse pages: {len(current_page)} + {len(next_page)} items'
                                )
                                i += 2  # Skip next page as it's been merged
                                continue

                        # If page is still too sparse and can't be merged, filter it out
                        if is_sparse and (len(current_page) < 2
                                          or page_text_length < 50):
                            logger.info(
                                f'Filtering out sparse page with '
                                f'{len(current_page)} items, {page_text_length} chars'
                            )
                            i += 1
                            continue

                        # Page is acceptable, keep it
                        filtered_pages.append(current_page)
                        i += 1

                    content_pages = filtered_pages

                    # If no content pages, create at least one empty slide
                    if not content_pages:
                        content_pages = [[]]

                    # Create slides for each content page
                    for page_idx, page_content in enumerate(content_pages):
                        # Add slide
                        slide_layout = prs.slide_layouts[
                            1]  # Title and content layout
                        slide = prs.slides.add_slide(slide_layout)

                        # Set title with adaptive font size (add page number if multiple pages)
                        if slide_data['title']:
                            title = slide.shapes.title
                            title_text = slide_data['title']
                            if len(content_pages) > 1:
                                title_text += f' ({page_idx + 1}/{len(content_pages)})'
                            title.text = title_text

                            # Adaptive title font size based on slide dimensions
                            title_font_size = max(
                                PptxInches(0.35), slide_width / 25)

                            for paragraph in title.text_frame.paragraphs:
                                paragraph.font.size = int(title_font_size)
                                paragraph.font.bold = True
                                paragraph.alignment = PP_ALIGN.CENTER

                        # Add content for this page
                        if page_content:
                            # Check if this page contains special content (table/image)
                            special_content = [
                                item for item in page_content
                                if item[0] in ['table', 'image']
                            ]
                            text_content = [
                                item for item in page_content
                                if item[0] not in ['table', 'image']
                            ]

                            # Handle text content
                            if text_content:
                                content_shape = slide.shapes.placeholders[1]
                                text_frame = content_shape.text_frame
                                text_frame.clear()

                                for i, (content_type, *content_data
                                        ) in enumerate(text_content):
                                    if content_type == 'text':
                                        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                                        p.text = content_data[0]
                                        p.font.size = int(base_font_size)
                                        p.alignment = PP_ALIGN.LEFT
                                    elif content_type == 'bullet':
                                        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                                        p.text = content_data[0]
                                        p.level = 0
                                        p.font.size = int(bullet_font_size)
                                        p.alignment = PP_ALIGN.LEFT
                                    elif content_type == 'heading':
                                        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                                        p.text = content_data[0]
                                        p.font.bold = True
                                        p.font.size = int(heading_font_size)
                                        p.alignment = PP_ALIGN.LEFT

                            # Handle special content (tables and images are processed separately)
                            for content_type, *content_data in special_content:
                                if content_type == 'table':
                                    # Add table to slide
                                    table_data = content_data[0]
                                    if table_data:
                                        # Create a new slide for the table
                                        table_slide_layout = prs.slide_layouts[
                                            5]  # Blank layout for better control
                                        table_slide = prs.slides.add_slide(
                                            table_slide_layout)

                                        # Get slide dimensions for adaptive sizing
                                        slide_width = prs.slide_width
                                        slide_height = prs.slide_height

                                        # Add title manually for better control
                                        if slide_data['title']:
                                            title_shape = table_slide.shapes.add_textbox(
                                                PptxInches(0.5),
                                                PptxInches(0.2),
                                                slide_width - PptxInches(1),
                                                PptxInches(1))
                                            title_frame = title_shape.text_frame
                                            title_frame.margin_top = PptxInches(
                                                0.1)
                                            title_frame.margin_bottom = PptxInches(
                                                0.1)
                                            title_para = title_frame.paragraphs[
                                                0]
                                            title_para.text = f"{slide_data['title']} - Table"
                                            title_para.alignment = PP_ALIGN.CENTER
                                            # Adaptive title font size based on slide width
                                            title_font_size = max(
                                                PptxInches(0.3),
                                                slide_width / 30)
                                            title_para.font.size = int(
                                                title_font_size)
                                            title_para.font.bold = True

                                        # Calculate table dimensions
                                        max_cols = max(
                                            len(row) for row in
                                            table_data) if table_data else 0
                                        num_rows = len(table_data)

                                        if max_cols > 0 and num_rows > 0:
                                            # Adaptive table positioning and sizing
                                            margin_left = slide_width * 0.05  # 5% margin
                                            margin_top = PptxInches(
                                                1.8
                                            ) if slide_data[
                                                'title'] else slide_height * 0.1

                                            table_width = slide_width - (
                                                    margin_left * 2)
                                            table_height = slide_height - margin_top - (
                                                    slide_height * 0.1)

                                            # Ensure minimum table size
                                            table_width = max(
                                                table_width, PptxInches(4))
                                            table_height = max(
                                                table_height, PptxInches(2))

                                            # Add table to slide with adaptive dimensions
                                            try:
                                                table = table_slide.shapes.add_table(
                                                    num_rows, max_cols,
                                                    margin_left, margin_top,
                                                    table_width,
                                                    table_height).table

                                                # Calculate adaptive font size based on table size and content
                                                # Base font size on table dimensions and number of cells
                                                cell_area = (table_width * table_height) / (num_rows * max_cols)
                                                base_font_size = max(PptxInches(0.12), min(PptxInches(0.25), cell_area / PptxInches(2)))

                                                # Fill table data with adaptive formatting
                                                for row_idx, row_data in enumerate(
                                                        table_data):
                                                    for col_idx in range(
                                                            max_cols):
                                                        try:
                                                            cell = table.cell(row_idx, col_idx)

                                                            # Get cell data or empty string if not enough columns
                                                            cell_data = str(row_data[col_idx] if col_idx < len(row_data) else '')

                                                            # Clear existing content first
                                                            cell.text = ''

                                                            # Add text to cell
                                                            if cell.text_frame.paragraphs:
                                                                paragraph = cell.text_frame.paragraphs[0]
                                                            else:
                                                                paragraph = cell.text_frame.add_paragraph()

                                                            paragraph.text = cell_data
                                                            paragraph.alignment = PP_ALIGN.CENTER

                                                            # Set font properties safely
                                                            if paragraph.font:
                                                                try:
                                                                    paragraph.font.size = int(base_font_size)
                                                                    # Make header row bold and slightly larger
                                                                    if row_idx == 0:  # First row is header
                                                                        paragraph.font.bold = True
                                                                        paragraph.font.size = int(base_font_size * 1.1)
                                                                except Exception as e:
                                                                    logger.warning(
                                                                        f'Error setting font size: {e}'
                                                                    )

                                                            # Set cell margins safely
                                                            try:
                                                                cell.margin_left = PptxInches(0.05)  # noqa
                                                                cell.margin_right = PptxInches(0.05)  # noqa
                                                                cell.margin_top = PptxInches(0.05)  # noqa
                                                                cell.margin_bottom = PptxInches(0.05)  # noqa
                                                            except Exception as e:
                                                                logger.warning(
                                                                    f'Error setting cell margins: {e}'
                                                                )

                                                        except Exception as cell_error:
                                                            logger.warning(
                                                                f'Error formatting table cell '
                                                                f'[{row_idx}][{col_idx}]: {cell_error}'
                                                            )
                                                            continue

                                            except Exception as table_error:
                                                logger.warning(
                                                    f'Error creating table: {table_error}'
                                                )
                                                # Fallback: create text-based table representation
                                                fallback_shape = table_slide.shapes.add_textbox(
                                                    margin_left, margin_top,
                                                    table_width, table_height)
                                                fallback_frame = fallback_shape.text_frame
                                                fallback_para = fallback_frame.paragraphs[0]
                                                fallback_text = '\n'.join([' | '.join(row) for row in table_data])
                                                fallback_para.text = f'Table Content:\n{fallback_text}'
                                elif content_type == 'image':
                                    img_path, alt_text = content_data[0], content_data[1] if len(content_data) > 1 else ''

                                    # Handle image path
                                    if not os.path.isabs(img_path):
                                        actual_img_path = resources_dir / img_path
                                        if not actual_img_path.exists():
                                            actual_img_path = markdown_folder / img_path
                                    else:
                                        actual_img_path = Path(img_path)

                                    if actual_img_path.exists():
                                        try:
                                            # Add image in new slide
                                            img_slide_layout = prs.slide_layouts[6]  # Blank layout
                                            img_slide = prs.slides.add_slide(img_slide_layout)

                                            # Get slide dimensions
                                            slide_width = prs.slide_width
                                            slide_height = prs.slide_height

                                            # Calculate available space for image (leaving margins)
                                            margin_top = PptxInches(1.5) if alt_text else PptxInches(0.5)
                                            margin_bottom = PptxInches(0.5)
                                            margin_left = PptxInches(0.5)
                                            margin_right = PptxInches(0.5)

                                            available_width = slide_width - margin_left - margin_right
                                            available_height = slide_height - margin_top - margin_bottom

                                            # Add image title at the top - always show title
                                            title_text = ''
                                            if alt_text and alt_text.strip():
                                                title_text = alt_text.strip()
                                            else:
                                                # If no alt_text, try to extract from image filename or use default
                                                img_name = Path(img_path).stem
                                                if img_name and not img_name.lower().startswith('image'):
                                                    title_text = f'Image: {img_name}'
                                                else:
                                                    title_text = 'Image'

                                            # Always add title
                                            if title_text:
                                                try:
                                                    title_shape = img_slide.shapes.add_textbox(
                                                        margin_left,
                                                        PptxInches(0.2),
                                                        available_width,
                                                        PptxInches(1))
                                                    title_frame = title_shape.text_frame
                                                    title_frame.margin_top = PptxInches(0.1)
                                                    title_frame.margin_bottom = PptxInches(0.1)
                                                    title_para = title_frame.paragraphs[0]
                                                    title_para.text = title_text
                                                    title_para.alignment = PP_ALIGN.CENTER
                                                    try:
                                                        title_para.font.size = int(PptxInches(0.25))
                                                        title_para.font.bold = True
                                                    except Exception as e:
                                                        logger.warning(
                                                            f'Error setting title font: {e}'
                                                        )
                                                        pass  # Skip font formatting if it fails
                                                    # Adjust margin for title
                                                    margin_top = PptxInches(1.5)
                                                except Exception as title_error:
                                                    logger.warning(
                                                        f'Error adding image title: {title_error}'
                                                    )

                                            # Get original image dimensions to calculate aspect ratio
                                            try:
                                                from PIL import Image as PILImage
                                                with PILImage.open(actual_img_path) as pil_img:
                                                    orig_width, orig_height = pil_img.size
                                                    aspect_ratio = orig_width / orig_height if orig_height > 0 else 4 / 3  # noqa
                                            except Exception as e:
                                                # Fallback aspect ratio if PIL is not available
                                                logger.warning(
                                                    f'Error loading image for aspect ratio: {e}'
                                                )
                                                aspect_ratio = 4 / 3

                                            # Recalculate available space after title
                                            available_height = slide_height - margin_top - margin_bottom

                                            # Calculate optimal image size to fit available space
                                            if available_width / available_height > aspect_ratio:
                                                # Height is the limiting factor
                                                img_height = available_height
                                                img_width = img_height * aspect_ratio
                                            else:
                                                # Width is the limiting factor
                                                img_width = available_width
                                                img_height = img_width / aspect_ratio

                                            # Center the image in available space
                                            img_left = margin_left + (available_width - img_width) / 2
                                            img_top = margin_top + (available_height - img_height) / 2

                                            # Add image with calculated dimensions
                                            try:
                                                img_slide.shapes.add_picture(
                                                    str(actual_img_path),
                                                    img_left,
                                                    img_top,
                                                    width=img_width,
                                                    height=img_height)
                                            except Exception as img_error:
                                                logger.warning(
                                                    f'Error adding image to slide: {img_error}'
                                                )
                                                # Add error text instead
                                                error_shape = img_slide.shapes.add_textbox(
                                                    img_left, img_top,
                                                    img_width, img_height)
                                                error_frame = error_shape.text_frame
                                                error_para = error_frame.paragraphs[0]
                                                error_para.text = f'Image loading failed: {img_path}'
                                                error_para.alignment = PP_ALIGN.CENTER

                                        except Exception as e:
                                            logger.warning(
                                                f'Failed to add image slide: {e}'
                                            )
                                            # Create simple error slide
                                            try:
                                                error_slide_layout = prs.slide_layouts[1]  # Title and content layout
                                                error_slide = prs.slides.add_slide(error_slide_layout)
                                                if error_slide.shapes.title:
                                                    error_slide.shapes.title.text = 'Image Loading Failed'
                                                if len(error_slide.shapes.placeholders) > 1:
                                                    content_shape = error_slide.shapes.placeholders[1]
                                                    if content_shape.text_frame:
                                                        p = content_shape.text_frame.paragraphs[0]
                                                        p.text = f'Unable to load image: {img_path}'
                                            except Exception as e:
                                                # Skip error slide creation if it also fails
                                                logger.warning(
                                                    f'Also failed to create error slide: {e}'
                                                )
                                    else:
                                        logger.warning(
                                            f'Image file not found: {actual_img_path}'
                                        )

                # If no meaningful content slides, create a summary slide
                if not slides_content:
                    try:
                        slide_layout = prs.slide_layouts[0]  # Title slide
                        slide = prs.slides.add_slide(slide_layout)
                        if slide.shapes.title:
                            slide.shapes.title.text = f'Document Content - {markdown_file.stem}'
                        if len(slide.shapes.placeholders) > 1:
                            subtitle = slide.shapes.placeholders[1]
                            if subtitle.text_frame:
                                subtitle.text = (
                                    f'Source file: {markdown_file.name}\n\nNote: The original document '
                                    f'may only contain titles without actual content.'
                                    f'\nPlease check the content structure of the original Markdown file.'
                                )
                    except Exception as default_slide_error:
                        logger.info(
                            f'Error creating default slide: {default_slide_error}'
                        )

                # Generate output PPT file path
                if single_file_output:
                    output_ppt_file = single_file_output
                else:
                    output_ppt_file = ppt_folder / f'{markdown_file.stem}.pptx'

                # Save PPT
                prs.save(output_ppt_file)
                converted_files.append(output_ppt_file)
                logger.info(f'✓ PPT file generated: {output_ppt_file}')

            except Exception as e:
                logger.error(f'✗ Failed to convert {markdown_file.name}: {e}')

        logger.info(
            f'\nConversion completed! Generated {len(converted_files)} PPT files.'
        )
        if single_file_output:
            return str(single_file_output.absolute())
        else:
            return str(ppt_folder.absolute())
